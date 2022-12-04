from flask import Flask, render_template, redirect, url_for, Blueprint, session, send_file
from GFX_APP import db, app
from GFX_APP.models import Project, NameKey, DefaultBugs, Socials, Slido, Locators, URLs
from pptx import Presentation
import os
from io import BytesIO, StringIO


SHORT_URL = 0
LONG_URL = 1
TWO_LINE_LEFT_SHORT = 4
TWO_LINE_LEFT_LONG = 5
TWO_LINE_RIGHT_SHORT = 6
TWO_LINE_RIGHT_LONG = 7
THREE_LINE_LEFT_SHORT = 8
THREE_LINE_LEFT_LONG = 9
THREE_LINE_RIGHT_SHORT = 10
THREE_LINE_RIGHT_LONG = 11
Q_A_BUG = 14
CHAT_BUG = 15
POLL_BUG = 16
SOCIAL_BUG = 17
SURVEY_BUG = 18
LOCATOR_BUG = 19
SLIDO_BUG = 20
RAISE_HAND = 21
ON_THE_PHONE_TWO_LINE = 25
ON_THE_PHONE_THREE_LINE = 26
LOCAL_GFX_PPT_LINK = "GFX_APP/static/pptx_template_file/gfx_29_2.pptx"
pres = Presentation(LOCAL_GFX_PPT_LINK)

# name_dictionary = {}
def create_url(url):
    url = url
    if len(url) <= 26:
        slide = pres.slides.add_slide(pres.slide_layouts[SHORT_URL])
    else:
        slide = pres.slides.add_slide(pres.slide_layouts[LONG_URL])
    url_placeholder = slide.placeholders[10]
    url_placeholder.text = url

def create_locator_bug(text):
    text = text
    slide = pres.slides.add_slide(pres.slide_layouts[LOCATOR_BUG])
    locator_placeholder = slide.placeholders[12]
    locator_placeholder.text = text

def create_slido_bug(event_code):
    event_code = event_code
    slide = pres.slides.add_slide(pres.slide_layouts[SLIDO_BUG])
    event_code_placeholder = slide.placeholders[12]
    join_placeholder = slide.placeholders[13]
    join_placeholder.text = "Join at slido.com"
    event_code_placeholder.text = event_code

def check_placeholders(template_num):
    slide = pres.slides.add_slide(pres.slide_layouts[template_num])
    for shape in slide.placeholders:
        print(f"{shape.placeholder_format.idx} {shape.name}")

def create_default_bug(bug):
    bug_passed = bug
    # print(bug_passed)
    bug_layout = pres.slide_layouts[bug_passed]
    bug_slide = pres.slides.add_slide(bug_layout)

def create_social_bug(text):
    text=text
    social_layout = pres.slide_layouts[SOCIAL_BUG]
    social_slide = pres.slides.add_slide(social_layout)
    text_placeholder = social_slide.placeholders[12]
    text_placeholder.text = text


def create_two_line_left_slide(n, t):
    name_passed = n
    title_passed = t
    if len(title_passed) <= 26 and len(name_passed) <= 27:
        two_line_left_short_layout = pres.slide_layouts[TWO_LINE_LEFT_SHORT]
        name_key_slide = pres.slides.add_slide(two_line_left_short_layout)
    else:
        two_line_left_long_layout = pres.slide_layouts[TWO_LINE_LEFT_LONG]
        name_key_slide = pres.slides.add_slide(two_line_left_long_layout)
    name_placeholder = name_key_slide.placeholders[10]
    title_placeholder = name_key_slide.placeholders[11]
    name_placeholder.text = name_passed
    title_placeholder.text = title_passed


def create_three_line_left_slide(n, t, c):
    name_passed = n
    title_passed = t
    company_passed = c
    if len(title_passed) <= 26 and len(name_passed) <= 27:
        three_line_left_short_layout = pres.slide_layouts[THREE_LINE_LEFT_SHORT]
        name_key_slide = pres.slides.add_slide(three_line_left_short_layout)
    else:
        three_line_left_long_layout = pres.slide_layouts[THREE_LINE_LEFT_LONG]
        name_key_slide = pres.slides.add_slide(three_line_left_long_layout)
    name_placeholder = name_key_slide.placeholders[10]
    title_placeholder = name_key_slide.placeholders[11]
    company_placeholder = name_key_slide.placeholders[12]
    name_placeholder.text = name_passed
    title_placeholder.text = title_passed
    company_placeholder.text = company_passed


def create_two_line_right_slide(n, t):
    name_passed = n
    title_passed = t
    if len(title_passed) <= 26 and len(name_passed) <= 27:
        two_line_right_short_layout = pres.slide_layouts[TWO_LINE_RIGHT_SHORT]
        name_key_slide = pres.slides.add_slide(two_line_right_short_layout)
    else:
        two_line_right_long_layout = pres.slide_layouts[TWO_LINE_RIGHT_LONG]
        name_key_slide = pres.slides.add_slide(two_line_right_long_layout)
    name_placeholder = name_key_slide.placeholders[10]
    title_placeholder = name_key_slide.placeholders[11]
    name_placeholder.text = name_passed
    title_placeholder.text = title_passed


def create_three_line_right_slide(n, t, c):
    name_passed = n
    title_passed = t
    company_passed = c
    if len(title_passed) <= 26 and len(name_passed) <= 27:
        three_line_right_short_layout = pres.slide_layouts[THREE_LINE_RIGHT_SHORT]
        name_key_slide = pres.slides.add_slide(three_line_right_short_layout)
    else:
        three_line_right_long_layout = pres.slide_layouts[THREE_LINE_RIGHT_LONG]
        name_key_slide = pres.slides.add_slide(three_line_right_long_layout)
    name_placeholder = name_key_slide.placeholders[10]
    title_placeholder = name_key_slide.placeholders[11]
    company_placeholder = name_key_slide.placeholders[12]
    name_placeholder.text = name_passed
    title_placeholder.text = title_passed
    company_placeholder.text = company_passed


def create_on_the_phone_slide(n, t, c):
    name_passed = n
    title_passed = t
    company_passed = c
    if company_passed == "":
        on_the_phone_layout = pres.slide_layouts[ON_THE_PHONE_TWO_LINE]
        on_the_phone_slide = pres.slides.add_slide(on_the_phone_layout)
    else:
        on_the_phone_layout = pres.slide_layouts[ON_THE_PHONE_THREE_LINE]
        on_the_phone_slide = pres.slides.add_slide(on_the_phone_layout)
        company_placeholder = on_the_phone_slide.placeholders[12]

    name_placeholder = on_the_phone_slide.placeholders[10]
    title_placeholder = on_the_phone_slide.placeholders[11]

    name_placeholder.text = name_passed
    title_placeholder.text = title_passed
    try:
        if company_passed != "":
            company_placeholder.text = company_passed
    except Exception as e:
        print(e)




view_all_blueprint = Blueprint('view_all', __name__, template_folder='templates/view_all')


@view_all_blueprint.route('/view_all', methods=['GET','POST'])
def view_all():
    cis = session['cis']
    # print(cis)
    project = Project.query.get(cis)
    try:
        name_keys = NameKey.query.filter_by(cis_id=cis).all()
        # print(name_keys)
    except Exception as e:
        print(e)
    bugs = DefaultBugs.query.filter_by(cis_id=cis).first()
    bug_list = []
    try:
        if bugs.q_a == 1:
            bug_list.append('Q & A')
        if bugs.chat == 1:
            bug_list.append(", Chat ")
        if bugs.poll == 1:
            bug_list.append(", Poll")
        if bugs.survey == 1:
            bug_list.append(", Survey")
        if bugs.raise_hand == 1:
            bug_list.append(", Raise Hand")
    except Exception as e:
        print(e)
        pass
    socials = Socials.query.filter_by(cis_id=cis).all()
    slidos = Slido.query.filter_by(cis_id=cis).all()
    locators = Locators.query.filter_by(cis_id=cis).all()
    urls = URLs.query.filter_by(cis_id=cis).all()

    return render_template('view_all.html', project=project, name_keys=name_keys, bugs_list=bug_list,
                           socials=socials, slidos=slidos, locators=locators, urls=urls)

@view_all_blueprint.route('/create_pptx', methods=['GET','POST'])
def create_pptx():
    cis = session['cis']
    project_name = session['project_name']
    slide_deck_name = f"{cis}_{project_name}_GFX.pptx"
    name_keys = NameKey.query.filter_by(cis_id=cis).all()

    for row in name_keys:
        name = row.name
        title = row.title
        company = row.company
        left_aligned = row.left_aligned
        right_aligned = row.right_aligned
        if company != "":
            if left_aligned:
                create_three_line_left_slide(name, title, company)
            if right_aligned:
                create_three_line_right_slide(name, title, company)
        else:
            if left_aligned:
                create_two_line_left_slide(name, title)
            if right_aligned:
                create_two_line_right_slide(name, title)
    for row in name_keys:
        name = row.name
        title = row.title
        company = row.company
        otp = row.otp
        if otp:
            create_on_the_phone_slide(name, title, company)



    bugs = DefaultBugs.query.filter_by(cis_id=cis).first()

    try:
        if bugs.q_a == 1:
            create_default_bug(Q_A_BUG)
        if bugs.chat == 1:
            create_default_bug(CHAT_BUG)
        if bugs.poll == 1:
            create_default_bug(POLL_BUG)
        if bugs.survey == 1:
            create_default_bug(SURVEY_BUG)
        if bugs.raise_hand == 1:
            create_default_bug(RAISE_HAND)
    except Exception as e:
        print(e)
        pass
    try:
        socials = Socials.query.filter_by(cis_id=cis).all()
        for social in socials:

            text = social.text
            create_social_bug(text)
    except Exception as e:
        print(e)

    slidos = Slido.query.filter_by(cis_id=cis).all()
    try:
        for bug in slidos:
            create_slido_bug(bug.event_code)
    except Exception as e:
        print(e)


    locators = Locators.query.filter_by(cis_id=cis).all()
    try:
        for locator in locators:
            create_locator_bug(locator.text)
    except Exception as e:
        print(e)

    urls = URLs.query.filter_by(cis_id=cis).all()
    try:
        for url in urls:
            create_url(url.url)
    except Exception as e:
        print(e)
    # check_placeholders(SHORT_URL)
    # check_placeholders(LONG_URL)

    file = BytesIO()
    pres.save(file)
    file.seek(0)

    return send_file(path_or_file=file, download_name=slide_deck_name, as_attachment=True,
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
